from __future__ import annotations

import asyncio
import logging
from pathlib import Path

import fitz  # PyMuPDF
import openpyxl
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)


_COMPRESS_SETTINGS: dict[str, dict] = {
    "light": {"garbage": 1, "deflate": True},
    "balanced": {"garbage": 3, "deflate": True, "deflate_images": True},
    "maximum": {"garbage": 4, "deflate": True, "deflate_images": True, "deflate_fonts": True, "clean": True},
}


async def compress_pdf(input_path: Path, output_path: Path, level: str = "balanced") -> None:
    settings = _COMPRESS_SETTINGS.get(level, _COMPRESS_SETTINGS["balanced"])

    def _run() -> None:
        doc = fitz.open(str(input_path))
        doc.save(str(output_path), **settings)
        doc.close()

    await asyncio.to_thread(_run)


async def get_pdf_info(input_path: Path) -> dict:
    def _run() -> dict:
        doc = fitz.open(str(input_path))
        count = len(doc)
        doc.close()
        return {"page_count": count}
    return await asyncio.to_thread(_run)


async def rotate_pdf(
    input_path: Path, output_path: Path, rotation: int, pages: list[int] | None = None
) -> None:
    def _run() -> None:
        doc = fitz.open(str(input_path))
        targets = pages if pages is not None else list(range(len(doc)))
        for i in targets:
            if 0 <= i < len(doc):
                page = doc[i]
                page.set_rotation((page.rotation + rotation) % 360)
        doc.save(str(output_path))
        doc.close()
    await asyncio.to_thread(_run)


async def reorder_pdf(input_path: Path, output_path: Path, order: list[int]) -> None:
    def _run() -> None:
        doc = fitz.open(str(input_path))
        if not order or max(order) >= len(doc) or min(order) < 0:
            raise ValueError("Invalid page order for this document.")
        doc.select(order)
        doc.save(str(output_path))
        doc.close()
    await asyncio.to_thread(_run)


async def watermark_pdf(
    input_path: Path,
    output_path: Path,
    text: str,
    opacity: float = 0.3,
    color: tuple[float, float, float] = (0.5, 0.5, 0.5),
) -> None:
    def _run() -> None:
        doc = fitz.open(str(input_path))
        try:
            for page in doc:
                rect = page.rect
                center = fitz.Point(rect.width / 2, rect.height / 2)
                fs = max(min(rect.width, rect.height) * 0.07, 8)
                font = fitz.Font("helv")
                text_len = font.text_length(text, fontsize=fs)
                pos = fitz.Point(center.x - text_len / 2, center.y + fs * 0.3)
                try:
                    tw = fitz.TextWriter(rect)
                    tw.append(pos, text, font=font, fontsize=fs)
                    tw.write_text(page, color=color, opacity=opacity, morph=(center, fitz.Matrix(-45)))
                except Exception:
                    pass  # skip watermark on pages too small to render text
            doc.save(str(output_path))
        finally:
            doc.close()
    await asyncio.to_thread(_run)


async def protect_pdf(input_path: Path, output_path: Path, password: str) -> None:
    def _run() -> None:
        doc = fitz.open(str(input_path))
        perm = (
            fitz.PDF_PERM_ACCESSIBILITY
            | fitz.PDF_PERM_PRINT
            | fitz.PDF_PERM_COPY
            | fitz.PDF_PERM_ANNOTATE
        )
        doc.save(
            str(output_path),
            encryption=fitz.PDF_ENCRYPT_AES_256,
            user_pw=password,
            owner_pw=password + "_owner",
            permissions=perm,
        )
        doc.close()
    await asyncio.to_thread(_run)


async def unlock_pdf(input_path: Path, output_path: Path, password: str) -> None:
    def _run() -> None:
        doc = fitz.open(str(input_path))
        try:
            if not doc.needs_pass:
                raise ValueError("This PDF is not password protected.")
            rc = doc.authenticate(password)
            if not rc:
                raise ValueError("Incorrect password — please check and try again.")
            doc.save(str(output_path), encryption=fitz.PDF_ENCRYPT_NONE)
        finally:
            doc.close()
    await asyncio.to_thread(_run)


async def pdf_to_word(input_path: Path, output_path: Path) -> None:
    """Convert PDF → DOCX using pdf2docx."""

    def _run() -> None:
        cv = Converter(str(input_path))
        cv.convert(str(output_path), start=0, end=None)
        cv.close()

    await asyncio.to_thread(_run)


async def pdf_to_excel(input_path: Path, output_path: Path) -> None:
    """
    Extract all tables from a PDF and write each page's tables
    as separate sheets in an Excel workbook.
    Falls back to raw text rows if no tables are detected.
    """

    def _run() -> None:
        doc = fitz.open(str(input_path))
        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # remove default empty sheet
        sheets_created = 0

        for page_num, page in enumerate(doc, start=1):
            tabs = page.find_tables()
            if tabs.tables:
                for t_idx, table in enumerate(tabs.tables, start=1):
                    rows = table.extract()
                    if not rows:
                        continue
                    sheet_name = f"P{page_num}-T{t_idx}"[:31]
                    ws = wb.create_sheet(title=sheet_name)
                    for row in rows:
                        ws.append([cell if cell is not None else "" for cell in row])
                    sheets_created += 1
            else:
                # No structured table — fall back to extracting lines of text
                text = page.get_text("text").strip()
                if text:
                    sheet_name = f"Page {page_num}"[:31]
                    ws = wb.create_sheet(title=sheet_name)
                    for line in text.splitlines():
                        ws.append([line])
                    sheets_created += 1

        if sheets_created == 0:
            ws = wb.create_sheet(title="Sheet1")
            ws.append(["No extractable content found in this PDF."])

        doc.close()
        wb.save(str(output_path))

    await asyncio.to_thread(_run)


# ---------------------------------------------------------------------------
# Repair PDF
# ---------------------------------------------------------------------------

async def repair_pdf(input_path: Path, output_path: Path) -> None:
    """Re-open and re-save with PyMuPDF — fixes most structural corruption."""

    def _run() -> None:
        doc = fitz.open(str(input_path))
        if len(doc) == 0:
            doc.close()
            raise ValueError("PDF has no pages and cannot be repaired.")
        doc.save(str(output_path), garbage=4, deflate=True, clean=True)
        doc.close()

    await asyncio.to_thread(_run)


# ---------------------------------------------------------------------------
# Page Numbers
# ---------------------------------------------------------------------------

_PAGE_NUMBER_POSITIONS = {
    "bottom-center": lambda w, h, tw, margin: ((w - tw) / 2, h - margin),
    "bottom-right":  lambda w, h, tw, margin: (w - tw - margin, h - margin),
    "bottom-left":   lambda w, h, tw, margin: (margin, h - margin),
    "top-center":    lambda w, h, tw, margin: ((w - tw) / 2, margin),
    "top-right":     lambda w, h, tw, margin: (w - tw - margin, margin),
    "top-left":      lambda w, h, tw, margin: (margin, margin),
}


async def add_page_numbers(
    input_path: Path,
    output_path: Path,
    position: str = "bottom-center",
    font_size: float = 11.0,
    start_number: int = 1,
    margin: float = 28.0,
) -> None:
    pos_fn = _PAGE_NUMBER_POSITIONS.get(position, _PAGE_NUMBER_POSITIONS["bottom-center"])

    def _run() -> None:
        doc = fitz.open(str(input_path))
        for i, page in enumerate(doc):
            label = str(start_number + i)
            # Estimate text width: helvetica is ~0.5× font_size per char
            tw = font_size * 0.5 * len(label)
            w, h = page.rect.width, page.rect.height
            x, y = pos_fn(w, h, tw, margin)
            page.insert_text(
                fitz.Point(x, y),
                label,
                fontname="helv",
                fontsize=font_size,
                color=(0.2, 0.2, 0.2),
            )
        doc.save(str(output_path))
        doc.close()

    await asyncio.to_thread(_run)


# ---------------------------------------------------------------------------
# PDF to Markdown
# ---------------------------------------------------------------------------

async def pdf_to_markdown(input_path: Path, output_path: Path) -> None:
    """Extract text from each page and write structured Markdown."""

    def _run() -> None:
        doc = fitz.open(str(input_path))
        lines: list[str] = [f"# {input_path.stem}\n"]
        for i, page in enumerate(doc, start=1):
            lines.append(f"\n## Page {i}\n")
            blocks = page.get_text("blocks")
            for b in blocks:
                text = b[4].strip()
                if not text:
                    continue
                # Heuristic: short lines in large font → heading
                fs = b[3] - b[1]  # block height proxy
                line_count = text.count("\n") + 1
                avg_h = fs / max(line_count, 1)
                if avg_h > 18 and len(text) < 120:
                    lines.append(f"\n### {text}\n")
                else:
                    lines.append(text + "\n")
        doc.close()
        output_path.write_text("\n".join(lines), encoding="utf-8")

    await asyncio.to_thread(_run)


# ---------------------------------------------------------------------------
# Crop PDF
# ---------------------------------------------------------------------------

_MM_TO_PT = 2.8346456692913385  # 1 mm = 2.835 pt


async def crop_pdf(
    input_path: Path,
    output_path: Path,
    top: float = 0.0,
    bottom: float = 0.0,
    left: float = 0.0,
    right: float = 0.0,
) -> None:
    """Crop each page by removing margins (in mm) from each side."""

    def _run() -> None:
        doc = fitz.open(str(input_path))
        for page in doc:
            mb = page.mediabox  # original full page rect (pts)
            new_rect = fitz.Rect(
                mb.x0 + left * _MM_TO_PT,
                mb.y0 + top * _MM_TO_PT,
                mb.x1 - right * _MM_TO_PT,
                mb.y1 - bottom * _MM_TO_PT,
            )
            if new_rect.is_empty or new_rect.is_infinite:
                raise ValueError("Crop margins are too large — nothing would remain on the page.")
            page.set_cropbox(new_rect)
        doc.save(str(output_path))
        doc.close()

    await asyncio.to_thread(_run)


# ---------------------------------------------------------------------------
# PDF to PowerPoint
# ---------------------------------------------------------------------------

async def pdf_to_pptx(input_path: Path, output_path: Path) -> None:
    """Render each PDF page as an image and embed it in a PPTX slide."""

    def _run() -> None:
        doc = fitz.open(str(input_path))
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # completely blank

        for page in doc:
            # Render at 150 DPI
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")

            slide = prs.slides.add_slide(blank_layout)
            import io
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        doc.close()
        prs.save(str(output_path))

    await asyncio.to_thread(_run)
